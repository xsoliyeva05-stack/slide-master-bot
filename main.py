import os
import asyncio
import requests
import google.generativeai as genai
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import Command
from aiogram.types import FSInputFile, InlineKeyboardMarkup, InlineKeyboardButton
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- KONFIGURATSIYA (Kalitlaringiz joylashtirildi) ---
TOKEN = "8461901986:AAHIQLMa1RckCqGCU71PJuJZCCnfKdWjYXk"
GEMINI_KEY = "AIzaSyBtUB1yq7lZqF29RPozUiIpj0DT9Rh5eU8"
PEXELS_KEY = "RCX0q9HsWJG4wb1HEvn7XjHUWteszFXWcjvwHipW69I9UZ5cGpzGLcSu"

genai.configure(api_key=GEMINI_KEY)
ai_model = genai.GenerativeModel('gemini-pro')
bot = Bot(token=TOKEN)
dp = Dispatcher()

# --- DIZAYN VA SHRIFTLAR FUNKSIYASI ---
def apply_advanced_style(slide, style):
    background = slide.background
    fill = background.fill
    fill.solid()
    
    # Standart sozlamalar
    font_name = 'Arial'
    
    if style == "dark":
        fill.foreground_color.rgb = RGBColor(20, 24, 35) # To'q kulrang/qora
        text_color = RGBColor(255, 255, 255) # Oq matn
        font_name = 'Verdana'
    elif style == "business":
        fill.foreground_color.rgb = RGBColor(7, 42, 108) # Professional ko'k
        text_color = RGBColor(255, 255, 255) # Oq matn
        font_name = 'Calibri'
    else: # Classic/White
        fill.foreground_color.rgb = RGBColor(255, 255, 255) # Sof oq
        text_color = RGBColor(30, 30, 30) # To'q kulrang matn
        font_name = 'Georgia'
        
    return text_color, font_name

# --- PEXELS'DAN RASM OLISH ---
def fetch_pexels_image(query):
    headers = {"Authorization": PEXELS_KEY}
    url = f"https://api.pexels.com/v1/search?query={query}&per_page=1"
    try:
        r = requests.get(url, headers=headers).json()
        if r.get('photos'):
            img_url = r['photos'][0]['src']['large']
            img_content = requests.get(img_url).content
            with open("temp_slide_img.jpg", "wb") as f:
                f.write(img_content)
            return "temp_slide_img.jpg"
    except:
        return None

# --- BOT INTERFEYSI ---
@dp.message(Command("start"))
async def welcome(message: types.Message):
    await message.answer("🌟 **Professional Slayd Botga xush kelibsiz!**\n\n"
                         "Men sizga AI yordamida matn yozib, rasmlar bilan bezatilgan "
                         "PowerPoint slaydlarini tayyorlab beraman.\n\n"
                         "📝 **Mavzu yuboring:** (Masalan: O'zbekiston turizmi)")

@dp.message(F.text)
async def ask_design(message: types.Message):
    topic = message.text
    buttons = [
        [InlineKeyboardButton(text="🌑 Dark Style (Zamonaviy)", callback_data=f"set_dark_{topic}")],
        [InlineKeyboardButton(text="🔵 Business Blue (Rasmiy)", callback_data=f"set_business_{topic}")],
        [InlineKeyboardButton(text="⚪ Classic White (Oddiy)", callback_data=f"set_classic_{topic}")]
    ]
    keyboard = InlineKeyboardMarkup(inline_keyboard=buttons)
    await message.answer(f"🎨 **'{topic}'** mavzusi uchun dizayn uslubini tanlang:", reply_markup=keyboard)

@dp.callback_query(F.data.startswith("set_"))
async def generate_presentation(call: types.CallbackQuery):
    _, style, topic = call.data.split("_")
    await call.message.edit_text("🔄 **AI ishlamoqda...**\nMatn tayyorlanmoqda va rasmlar yig'ilmoqda. Iltimos, kuting.")
    
    try:
        # Gemini AI dan o'zbekcha kontent so'rash
        prompt = (f"Write 5 educational slides about '{topic}' in Uzbek language. "
                  "Structure: Start each slide with 'S:', then title. Then body text. "
                  "At the end of each slide, write 'R:' followed by one specific English keyword for a photo.")
        
        response = ai_model.generate_content(prompt)
        prs = Presentation()
        
        slides_content = response.text.split("S:")[1:]
        for content in slides_content:
            lines = content.strip().split('\n')
            title_text = lines[0]
            
            # Kalit so'zni va matnni ajratish
            img_keyword = "presentation"
            clean_body = []
            for line in lines[1:]:
                if "R:" in line:
                    img_keyword = line.replace("R:", "").strip()
                else:
                    clean_body.append(line)
            
            # Slayd qo'shish
            slide_layout = prs.slide_layouts[1] # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Dizayn va Shriftlarni qo'llash
            t_color, f_name = apply_advanced_style(slide, style)
            
            # 1. Sarlavha dizayni
            title_shape = slide.shapes.title
            title_shape.text = title_text
            title_frame = title_shape.text_frame.paragraphs[0]
            title_frame.font.bold = True
            title_frame.font.size = Pt(32)
            title_frame.font.name = f_name
            title_frame.font.color.rgb = t_color
            
            # 2. Asosiy matn dizayni
            body_shape = slide.placeholders[1]
            body_shape.text = "\n".join(clean_body)
            for para in body_shape.text_frame.paragraphs:
                para.font.size = Pt(18)
                para.font.name = f_name
                para.font.color.rgb = t_color
                para.alignment = PP_ALIGN.LEFT
            
            # 3. Pexels rasm dizayni
            image_path = fetch_pexels_image(img_keyword)
            if image_path:
                # Rasmni o'ng tomonga chiroyli joylashtirish
                slide.shapes.add_picture(image_path, Inches(6.5), Inches(1.5), Inches(3), Inches(3.5))

        # Saqlash va yuborish
        file_name = f"{topic.replace(' ', '_')}.pptx"
        prs.save(file_name)
        
        doc = FSInputFile(file_name)
        await bot.send_document(call.message.chat.id, doc, caption=f"✅ Slayd tayyor!\n📌 Mavzu: {topic}\n🎨 Uslub: {style}")
        
        # Tozalash
        os.remove(file_name)
        if os.path.exists("temp_slide_img.jpg"): os.remove("temp_slide_img.jpg")
        
    except Exception as e:
        await call.message.answer(f"❌ Xatolik yuz berdi: {str(e)}")
    finally:
        await call.message.delete()

async def main():
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
