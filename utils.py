from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO

def create_pptx(topic, slides_data, full_name):
    prs = Presentation()
    
    # 1-SLAYD: TITUL (DIZAYN)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = topic.upper()
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) # To'q ko'k
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    subtitle.text = f"Tayyorladi: {full_name}\nSlayd Master AI orqali yaratildi"
    
    # ASOSIY SLAYDLAR
    for info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        # Sarlavha dizayni
        title_shape = slide.shapes.title
        title_shape.text = info['title']
        for p in title_shape.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(0, 102, 204) # Moviy
            p.font.bold = True
            
        # Matn dizayni
        body_shape = slide.placeholders[1]
        body_shape.text = info['content']
        for p in body_shape.text_frame.paragraphs:
            p.font.size = Pt(18)
            
    file_stream = BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    return file_stream
  
